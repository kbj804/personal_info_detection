'''
pip install pdfminer.six  # https://pdfminersix.readthedocs.io/en/latest/tutorial/extract_pages.html
pip install 'olefile'  or '-U olefile'
pip install python-pptx
pip install python-docx
pip install openpyxl
'''
import os
from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.converter import HTMLConverter
from pdfminer.converter import TextConverter
from pdfminer.layout import LAParams
from pdfminer.pdfpage import PDFPage
from io import StringIO
import olefile
from pptx import Presentation
from docx import Document
import csv
import openpyxl
from pdfminer.high_level import extract_pages
from pdfminer.layout import LTTextContainer

class loadFileManager:
    # 파일 이름, 확장자 분리
    def __init__(self, path):
        # dir_path = os.getcwd()
        self.path = path
        basename = os.path.basename(self.path)
        self.name, self.dotext = os.path.splitext(basename)
        self.ext = self.dotext.replace(".",'',1)

        self.data = ''

        # 읽을 수 있는 파일인지 검사하고 읽을 수 있으면 data입력
        if self.check_ext():
            self.data = self.read_file()

        else:
            print("여기론 안와")
            pass
    
    # 확장자에 맞는 read 함수로 매핑
    def read_file(self):
        result = self.read_function[self.ext](self)
        return result

    # 읽을 수 있는 확장자인가 검사
    def check_ext(self):
        try:
            if self.read_function[self.ext]:
                return True
            else:
                print("@ # @ # {} For Debugging... @ # @ #".format(self.ext))
                return False

        except Exception as e:
            print("####ERROR#### {0} dose not exist in ext Dictionary".format(e))

    # 확장자 별 데이터 오픈
    # pdf, hwp, ppt, docx, ....

    def read_pdf(self):
        # using pdfminer.six  / No pdfminer
        result = {}
        for i, page_layout in enumerate(extract_pages(self.path)):
            page_contents=''
            for element in page_layout:
                if isinstance(element, LTTextContainer):
                    page_contents += element.get_text()
            result[i] = page_contents
        
        return result
        

        # pdf 추출하는 다른 방법 (pdfminer)
        '''rsrcmgr = PDFResourceManager()
        retstr = StringIO()
        codec = 'utf-8'
        laparams = LAParams()
        # f = open('./out.html', 'wb')
        # device = HTMLConverter(rsrcmgr, f, codec=codec, laparams=laparams)
        device = TextConverter(rsrcmgr, retstr, codec=codec, laparams=laparams)
        
        with open(self.path, 'rb') as fp:
            interpreter = PDFPageInterpreter(rsrcmgr, device)
            # password = ""
            maxpages = 0 #is for all
            # caching = True
            pagenos=set()
            
            for page in PDFPage.get_pages(fp, pagenos, maxpages=maxpages, check_extractable=True):
                interpreter.process_page(page)
            
            # 이건 전체 페이지 전부다 리턴할때 쓰는거
            str = retstr.getvalue()

            fp.close()
        
        device.close()
        retstr.close()
        f.close()
    
        return str
        '''

    def read_hwp(self):
        f = olefile.OleFileIO(self.path)
        encoded_text = f.openstream('PrvText').read()
        decoded_text = encoded_text.decode('UTF-16')
        
        return decoded_text

    def read_pptx(self):
        pptxdoc = Presentation(self.path)
        # text_runs will be populated with a list of strings,
        # one for each text run in presentation
        text_runs = []
        for slide in pptxdoc.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    # print("text")
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        # 1. run -> 모든 내용 제목 까지 전부다
                        for run in paragraph.runs:
                            text_runs.append(run.text)

                        # 2. paragraph -> 페이지별로 구분됨
                        # text_runs.append(paragraph.text)
                
                elif shape.has_table:
                    # print("table")
                    tb1 = shape.table
                    row_count = len(tb1.rows)
                    col_count = len(tb1.columns)
                    for r in range(0, row_count):
                        for c in range(0, col_count):
                            cell = tb1.cell(r,c)
                            paragraphs = cell.text_frame.paragraphs 
                            for paragraph in paragraphs:
                                for run in paragraph.runs:
                                    text_runs.append(run.text)
                else:
                    # print("etc")
                    pass
 
        return '\n'.join(text_runs)

    def read_docx(self):
        # with open(self.path, 'rb') as f:
        #     source_stream = StringIO(f.read())
        document = Document(self.path)
        fullText = []
        for para in document.paragraphs:
            fullText.append(para.text)
        # source_stream.close()

        return '\n'.join(fullText)

    # csv 코덱문제 해결하고 엑셀읽기 하고 html읽기 하면됨 
    def read_csv(self):
        read_list=[]
        with open(self.path, 'r', encoding='UTF8') as file:
            file_read = csv.reader(file)
            for line in file_read:
                read_list.append(line)

            # 2차원 리스트 -> 1차원으로
            answer = sum(read_list,[])
        
        return '\n'.join(answer)

    # 리스트 값 안에 숫자형도 있어서 \n . join 사용이 안됨
    # 리스트가 3차원 까지 가서 코드가 다소 복잡함
    # 마지막 리스트가 문자열로 치환이 안되어 추수 수정 필요
    def read_xlsx(self):
        workbook = openpyxl.load_workbook(self.path, data_only=True)
        # Sheet 목록
        sheet_list = workbook.sheetnames
        all_sheet_value=[]
        
        # Sheet 별 탐색
        for sheet in sheet_list:
            all_values = []
            workSheet = workbook[sheet]
            for row in workSheet.rows:
                row_value =[]
                for cell in row:
                    if cell.value is None:
                        pass
                    else:
                        row_value.append(cell.value)
                all_values.append(row_value)
            # 차원축소1
            answer = sum(all_values,[])
            all_sheet_value.append(answer)
        # 차원축소2
        answer2 = sum(all_sheet_value, [])

        # 리스트 강제로 str로 함 ㅠ
        return str(answer2)
        

    def read_txt(self):
        with open(self.path, 'r', encoding='UTF8') as f:
            data = f.read()
            return data
    
    def read_html(self):
        pass

    read_function = {
        'pdf': read_pdf,
        'hwp': read_hwp,
        'pptx': read_pptx,
        'docx': read_docx,
        'csv': read_csv,
        'xlsx': read_xlsx,
        'txt': read_txt,
        'html': read_html
    }
    
    

# a = loadFileManager('hwp_sample.hwp')
# print(a.data)
